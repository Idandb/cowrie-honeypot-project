#!/usr/bin/env python3
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


def parse_markdown_slides(md_text: str):
    lines = md_text.splitlines()

    # Skip YAML front matter if present
    idx = 0
    if idx < len(lines) and lines[idx].strip() == '---':
        idx += 1
        while idx < len(lines) and lines[idx].strip() != '---':
            idx += 1
        if idx < len(lines) and lines[idx].strip() == '---':
            idx += 1

    slides = []
    current = None

    heading_re = re.compile(r'^##\s+(.*)')

    while idx < len(lines):
        line = lines[idx]
        idx += 1

        m = heading_re.match(line)
        if m:
            # Start new slide
            title_raw = m.group(1).strip()
            # If heading like "שקופית N: כותרת", take text after ':'
            title_parts = [p.strip() for p in title_raw.split(':', 1)]
            title = title_parts[1] if len(title_parts) == 2 else title_raw
            if current:
                slides.append(current)
            current = {"title": title, "bullets": [], "paras": []}
            continue

        if current is None:
            # ignore any text before first heading
            continue

        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith('- '):
            bullet_text = stripped[2:].strip()
            current["bullets"].append(bullet_text)
        else:
            current["paras"].append(stripped)

    if current:
        slides.append(current)

    return slides


def add_slide(prs: Presentation, title: str, bullets: list[str], paras: list[str]):
    # Use Title and Content layout (usually index 1)
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # Title
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = title
        for p in title_shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.RIGHT
            if p.runs:
                for r in p.runs:
                    r.font.size = Pt(36)
                    r.font.name = 'Arial'

    # Content
    body = None
    for placeholder in slide.placeholders:
        if placeholder.has_text_frame and placeholder != title_shape:
            body = placeholder
            break
    if body is None:
        return

    tf = body.text_frame
    tf.clear()

    has_any = False
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if has_any else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.alignment = PP_ALIGN.RIGHT
        for r in p.runs:
            r.font.size = Pt(20)
            r.font.name = 'Arial'
        has_any = True

    for para in paras:
        p = tf.add_paragraph() if has_any else tf.paragraphs[0]
        p.text = para
        p.level = 0
        p.alignment = PP_ALIGN.RIGHT
        for r in p.runs:
            r.font.size = Pt(20)
            r.font.name = 'Arial'
        has_any = True


def convert(md_path: Path, pptx_path: Path):
    md_text = md_path.read_text(encoding='utf-8')
    slides = parse_markdown_slides(md_text)

    prs = Presentation()
    # Remove default first slide if any templates assume one; we'll just start from scratch
    # (python-pptx starts empty by default)

    for s in slides:
        add_slide(prs, s['title'], s['bullets'], s['paras'])

    prs.save(str(pptx_path))


if __name__ == '__main__':
    md_file = Path('/workspace/slides.md')
    out_file = Path('/workspace/slides.pptx')
    if not md_file.exists():
        raise SystemExit(f"Markdown file not found: {md_file}")
    out_file.parent.mkdir(parents=True, exist_ok=True)
    convert(md_file, out_file)
    print(f"Wrote {out_file}")

